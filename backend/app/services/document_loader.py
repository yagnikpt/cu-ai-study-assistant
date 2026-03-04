"""Unified document loader.

Supports PDF, DOCX, and PPTX formats. Auto-detects file type from filename
extension and dispatches to the appropriate loader.

Each loader extracts:
  - Page-by-page text content (with heading detection)
  - Embedded images with metadata (raw bytes, page number, mime type)

Usage:
    from app.services.document_loader import load_document

    parsed = load_document("lecture.pdf", pdf_bytes)
    print(parsed.pages[0].text)
    for img in parsed.images:
        print(img.mime_type, img.page_number, len(img.data))
"""

import hashlib
import io
import logging
from dataclasses import dataclass, field
from pathlib import PurePosixPath

logger = logging.getLogger(__name__)

# Supported file extensions (lowercase, with leading dot)
SUPPORTED_EXTENSIONS: set[str] = {".pdf", ".docx", ".pptx"}


@dataclass
class PageContent:
    """Extracted content from a single document page/slide."""

    page_number: int  # 1-indexed
    text: str
    headings: list[str] = field(default_factory=list)


@dataclass
class ImageData:
    """An image extracted from a document, with metadata."""

    data: bytes  # raw image bytes
    mime_type: str  # e.g. "image/png", "image/jpeg"
    page_number: int | None = None  # page the image was found on (None if unknown)


@dataclass
class ParsedDocument:
    """Result of parsing an entire document."""

    pages: list[PageContent]
    page_count: int
    metadata: dict
    images: list[ImageData] = field(default_factory=list)


def load_document(filename: str, source: bytes) -> ParsedDocument:
    """Parse a document and extract text + images.

    Auto-detects the file format from the filename extension and dispatches
    to the appropriate loader.

    Args:
        filename: Original filename (used to detect format, e.g. "notes.pdf").
        source: Raw file bytes (read entirely into memory).

    Returns:
        ParsedDocument with page-by-page content, metadata, and images.

    Raises:
        ValueError: If the file type is unsupported or the document is invalid.
    """
    ext = PurePosixPath(filename).suffix.lower()

    if ext == ".pdf":
        return _load_pdf(source)
    elif ext == ".docx":
        return _load_docx(source)
    elif ext == ".pptx":
        return _load_pptx(source)
    else:
        raise ValueError(
            f"Unsupported file type: '{ext}'. "
            f"Supported formats: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )


# ── PDF Loader (PyMuPDF) ───────────────────────────────


def _load_pdf(source: bytes) -> ParsedDocument:
    """Parse a PDF using PyMuPDF. Extracts text and images."""
    import pymupdf

    try:
        doc = pymupdf.open(stream=source, filetype="pdf")
    except Exception as e:
        raise ValueError(f"Failed to open PDF: {e}") from e

    if doc.is_encrypted:
        doc.close()
        raise ValueError("PDF is encrypted and cannot be parsed")

    # Document-level metadata
    pdf_metadata = doc.metadata or {}
    metadata = {
        "title": pdf_metadata.get("title", ""),
        "author": pdf_metadata.get("author", ""),
        "subject": pdf_metadata.get("subject", ""),
        "keywords": pdf_metadata.get("keywords", ""),
        "page_count": doc.page_count,
        "format": "pdf",
    }

    pages: list[PageContent] = []
    images: list[ImageData] = []
    seen_xrefs: set[int] = set()  # deduplicate images shared across pages

    for page_idx in range(doc.page_count):
        page = doc[page_idx]
        page_number = page_idx + 1

        # ── Text extraction ──
        page_dict: dict = page.get_text("dict", flags=pymupdf.TEXT_PRESERVE_WHITESPACE)  # type: ignore[assignment]
        blocks: list[dict] = page_dict.get("blocks", [])

        page_text_parts: list[str] = []
        headings: list[str] = []

        for block in blocks:
            if block.get("type") != 0:  # skip non-text blocks
                continue

            block_text_lines: list[str] = []
            for line in block.get("lines", []):
                line_text = ""
                max_font_size = 0.0

                for span in line.get("spans", []):
                    line_text += span.get("text", "")
                    max_font_size = max(max_font_size, span.get("size", 0))

                line_text = line_text.strip()
                if not line_text:
                    continue

                block_text_lines.append(line_text)

                # Heuristic: lines with font size > 14pt are likely headings
                if max_font_size > 14 and len(line_text) < 200:
                    headings.append(line_text)

            block_text = "\n".join(block_text_lines)
            if block_text.strip():
                page_text_parts.append(block_text)

        page_text = "\n\n".join(page_text_parts)

        pages.append(
            PageContent(
                page_number=page_number,
                text=page_text,
                headings=headings,
            )
        )

        # ── Image extraction ──
        for img_info in page.get_images(full=True):
            xref = img_info[0]
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)

            try:
                extracted = doc.extract_image(xref)
                if extracted and extracted.get("image"):
                    img_bytes: bytes = extracted["image"]
                    # Skip tiny images (likely icons/bullets, <1KB)
                    if len(img_bytes) < 1024:
                        continue
                    # PyMuPDF returns "ext" key with format like "png", "jpeg"
                    img_ext = extracted.get("ext", "png")
                    mime_type = f"image/{img_ext}" if img_ext != "jpg" else "image/jpeg"
                    images.append(
                        ImageData(
                            data=img_bytes,
                            mime_type=mime_type,
                            page_number=page_number,
                        )
                    )
            except Exception:
                logger.debug(f"Failed to extract image xref={xref}, skipping")
                continue

    doc.close()

    logger.info(f"Parsed PDF: {len(pages)} pages, {len(images)} images extracted")

    return ParsedDocument(
        pages=pages,
        page_count=len(pages),
        metadata=metadata,
        images=images,
    )


# ── DOCX Loader (python-docx) ──────────────────────────


def _load_docx(source: bytes) -> ParsedDocument:
    """Parse a DOCX using python-docx. Extracts text and images.

    DOCX is a flow format without real page numbers. We treat the entire
    document as page 1 for all paragraphs. Headings are detected from
    paragraph style names (e.g. 'Heading 1', 'Heading 2', 'Title').
    """
    from docx import Document as DocxDocument
    from docx.opc.constants import RELATIONSHIP_TYPE as RT

    try:
        doc = DocxDocument(io.BytesIO(source))
    except Exception as e:
        raise ValueError(f"Failed to open DOCX: {e}") from e

    # Document-level metadata from core properties
    props = doc.core_properties
    metadata: dict[str, str | int] = {
        "title": props.title or "",
        "author": props.author or "",
        "subject": props.subject or "",
        "keywords": props.keywords or "",
        "format": "docx",
    }

    # ── Text extraction ──
    # Heading style names that indicate a heading
    heading_style_prefixes = ("heading", "title", "subtitle")

    paragraphs_text: list[str] = []
    headings: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        paragraphs_text.append(text)

        # Detect headings by style name
        style_name = (para.style.name or "").lower()
        if any(style_name.startswith(prefix) for prefix in heading_style_prefixes):
            headings.append(text)

    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_texts = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    row_texts.append(cell_text)
            if row_texts:
                paragraphs_text.append(" | ".join(row_texts))

    full_text = "\n\n".join(paragraphs_text)

    pages = [
        PageContent(
            page_number=1,
            text=full_text,
            headings=headings,
        )
    ]

    # ── Image extraction ──
    images: list[ImageData] = []
    seen_hashes: set[str] = set()  # deduplicate images referenced multiple times

    for rel in doc.part.rels.values():
        if rel.reltype == RT.IMAGE:
            try:
                img_bytes: bytes = rel.target_part.blob
                # Skip tiny images (<1KB)
                if len(img_bytes) < 1024:
                    continue
                # Deduplicate by content hash
                img_hash = hashlib.sha256(img_bytes).hexdigest()
                if img_hash in seen_hashes:
                    continue
                seen_hashes.add(img_hash)
                # Determine mime type from content_type or filename
                content_type = getattr(rel.target_part, "content_type", "")
                if content_type and content_type.startswith("image/"):
                    mime_type = content_type
                else:
                    mime_type = "image/png"
                images.append(
                    ImageData(
                        data=img_bytes,
                        mime_type=mime_type,
                        page_number=1,  # DOCX is treated as single page
                    )
                )
            except Exception:
                logger.debug(f"Failed to extract DOCX image rel={rel.rId}, skipping")
                continue

    metadata["page_count"] = 1

    logger.info(
        f"Parsed DOCX: {len(paragraphs_text)} paragraphs, "
        f"{len(headings)} headings, {len(images)} images"
    )

    return ParsedDocument(
        pages=pages,
        page_count=1,
        metadata=metadata,
        images=images,
    )


# ── PPTX Loader (python-pptx) ──────────────────────────


def _load_pptx(source: bytes) -> ParsedDocument:
    """Parse a PPTX using python-pptx. Extracts text and images.

    Each slide is treated as a page. Title placeholder text is detected
    as headings. Images are extracted from picture shapes.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        prs = Presentation(io.BytesIO(source))
    except Exception as e:
        raise ValueError(f"Failed to open PPTX: {e}") from e

    # Presentation-level metadata (pptx core_properties)
    props = prs.core_properties
    metadata: dict[str, str | int] = {
        "title": props.title or "",
        "author": props.author or "",
        "subject": props.subject or "",
        "keywords": props.keywords or "",
        "format": "pptx",
    }

    pages: list[PageContent] = []
    images: list[ImageData] = []
    seen_hashes: set[str] = set()  # deduplicate images reused across slides

    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1
        text_parts: list[str] = []
        headings: list[str] = []

        for shape in slide.shapes:
            # ── Image extraction ──
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes: bytes = shape.image.blob
                    # Skip tiny images (<1KB)
                    if len(img_bytes) < 1024:
                        continue
                    # Deduplicate by content hash
                    img_hash = hashlib.sha256(img_bytes).hexdigest()
                    if img_hash in seen_hashes:
                        continue
                    seen_hashes.add(img_hash)
                    # python-pptx Image has content_type property
                    content_type = getattr(shape.image, "content_type", "image/png")
                    images.append(
                        ImageData(
                            data=img_bytes,
                            mime_type=content_type or "image/png",
                            page_number=slide_number,
                        )
                    )
                except Exception:
                    logger.debug(
                        f"Failed to extract PPTX image on slide {slide_number}, skipping"
                    )
                    continue

            # ── Text extraction ──
            if shape.has_text_frame:
                shape_text_parts: list[str] = []
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        shape_text_parts.append(para_text)

                shape_text = "\n".join(shape_text_parts)
                if shape_text.strip():
                    text_parts.append(shape_text)

                    # Title/subtitle placeholders are headings
                    if shape.is_placeholder:
                        ph_type = shape.placeholder_format.type
                        # PP_PLACEHOLDER: TITLE=1, CENTER_TITLE=3, SUBTITLE=4, VERTICAL_TITLE=5
                        if ph_type is not None and ph_type.value in (1, 3, 4, 5):
                            for line in shape_text_parts:
                                headings.append(line)

            # ── Table extraction ──
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        text_parts.append(" | ".join(row_texts))

        slide_text = "\n\n".join(text_parts)

        pages.append(
            PageContent(
                page_number=slide_number,
                text=slide_text,
                headings=headings,
            )
        )

    metadata["page_count"] = len(pages)

    logger.info(f"Parsed PPTX: {len(pages)} slides, {len(images)} images extracted")

    return ParsedDocument(
        pages=pages,
        page_count=len(pages),
        metadata=metadata,
        images=images,
    )
